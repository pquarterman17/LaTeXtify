"""Generate the privacy-fixture family: files that deliberately leak.

One script for the whole family (rather than the usual one-script-per-fixture)
because these fixtures exist to exercise a single subsystem,
:mod:`latextify.privacy`, and each one is only interesting relative to the
others.

Every fixture is built to carry a KNOWN set of leaks, listed in
``leaky_files.truth.json`` beside them, so tests assert against ground truth
rather than against whatever the code happens to report today.

The office fixtures are built by the real libraries (python-pptx, openpyxl) and
then post-processed at the ZIP level, because the leaks that matter most --
a hidden slide, an embedded chart workbook, a pivot cache -- are exactly the
things the friendly APIs will not let you create.

Run with: uv run python tests/fixtures/make_leaky_files.py
"""

from __future__ import annotations

import json
import re
import shutil
import zipfile
from pathlib import Path

FIXTURE_DIR = Path(__file__).parent
TRUTH_PATH = FIXTURE_DIR / "leaky_files.truth.json"

AUTHOR = "Dr Testy McTestface"
COMPANY = "Institute of Very Secret Research"


# ── zip post-processing helper ───────────────────────────────────────────


def rewrite_zip(
    path: Path,
    *,
    replacements: dict[str, bytes] | None = None,
    additions: dict[str, bytes] | None = None,
) -> None:
    """Rebuild ``path`` applying member replacements/additions."""
    replacements = replacements or {}
    additions = additions or {}
    tmp = path.with_suffix(path.suffix + ".tmp")
    with zipfile.ZipFile(path) as zin, zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = replacements.get(item.filename, zin.read(item.filename))
            zout.writestr(item.filename, data)
        for name, data in additions.items():
            zout.writestr(name, data)
    shutil.move(str(tmp), str(path))


def set_app_property(xml: bytes, name: str, value: str) -> bytes:
    """Set a docProps/app.xml property, replacing any existing element.

    Office writes empty placeholder elements (``<Company></Company>``), so
    appending a populated duplicate leaves the placeholder in front of it --
    which is exactly the shape of malformed file that hid this leak the first
    time round. Replace in place, or insert if genuinely absent.
    """
    element = f"<{name}>{value}</{name}>".encode()
    pattern = re.compile(rf"<{name}\s*/>|<{name}>.*?</{name}>".encode(), re.DOTALL)
    if pattern.search(xml):
        return pattern.sub(element, xml, count=1)
    return xml.replace(b"</Properties>", element + b"</Properties>")


# ── PowerPoint ───────────────────────────────────────────────────────────


def build_deck() -> dict[str, object]:
    from pptx import Presentation
    from pptx.util import Emu, Inches

    dest = FIXTURE_DIR / "leaky_deck.pptx"
    prs = Presentation()
    blank = prs.slide_layouts[6]
    title_layout = prs.slide_layouts[5]

    visible = prs.slides.add_slide(title_layout)
    visible.shapes.title.text = "Perfectly Ordinary Results"
    visible.notes_slide.notes_text_frame.text = (
        "Do NOT mention that the third replicate failed."
    )

    second = prs.slides.add_slide(blank)
    box = second.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Visible content"
    # Dragged off the left edge: never renders, still in the file.
    off = second.shapes.add_textbox(Emu(-4000000), Inches(1), Inches(3), Inches(1))
    off.text_frame.text = "Reviewer 2 is wrong and I will die on this hill"

    hidden = prs.slides.add_slide(blank)
    hbox = hidden.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    hbox.text_frame.text = "Unblinded patient identifiers"

    prs.core_properties.author = AUTHOR
    prs.core_properties.last_modified_by = "reviewer-laptop"
    prs.core_properties.title = "INTERNAL DRAFT - do not circulate"
    prs.save(dest)

    # Post-process: mark slide 3 hidden and embed a chart workbook.
    with zipfile.ZipFile(dest) as zin:
        names = zin.namelist()
        slide3 = sorted(n for n in names if n.startswith("ppt/slides/slide"))[2]
        slide3_xml = zin.read(slide3)
        app_xml = zin.read("docProps/app.xml")
        ct_xml = zin.read("[Content_Types].xml")

    hidden_xml = slide3_xml.replace(b"<p:sld ", b'<p:sld show="0" ', 1)
    assert hidden_xml != slide3_xml, "failed to mark slide 3 hidden"

    app_xml = set_app_property(app_xml, "Company", COMPANY)
    app_xml = set_app_property(app_xml, "TotalTime", "4242")
    ct_xml = ct_xml.replace(
        b"<Types ",
        b'<Types ',
    ).replace(
        b"</Types>",
        b'<Default Extension="xlsx" ContentType="application/vnd.openxmlformats-'
        b'officedocument.spreadsheetml.sheet"/></Types>',
    )

    workbook_bytes = _tiny_xlsx_bytes()
    rewrite_zip(
        dest,
        replacements={slide3: hidden_xml, "docProps/app.xml": app_xml,
                      "[Content_Types].xml": ct_xml},
        additions={"ppt/embeddings/Microsoft_Excel_Worksheet1.xlsx": workbook_bytes},
    )

    return {
        "path": dest.name,
        "author": AUTHOR,
        "company": COMPANY,
        "hidden_slides": 1,
        "notes_slides": 1,
        "embedded_workbooks": 1,
        "offcanvas_shapes": 1,
    }


def _tiny_xlsx_bytes() -> bytes:
    """A real (small) xlsx, used as PowerPoint's embedded chart workbook."""
    from openpyxl import Workbook

    tmp = FIXTURE_DIR / "_embed_tmp.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "SourceData"
    ws.append(["patient", "dose_mg", "outcome"])
    ws.append(["P-00417", 250, "adverse"])
    wb.save(tmp)
    data = tmp.read_bytes()
    tmp.unlink()
    return data


# ── Excel ────────────────────────────────────────────────────────────────


def build_workbook() -> dict[str, object]:
    from openpyxl import Workbook

    dest = FIXTURE_DIR / "leaky_book.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "Summary"
    ws.append(["region", "total"])
    ws.append(["North", 42])

    raw = wb.create_sheet("RawSalaries")
    raw.append(["name", "salary"])
    raw.append(["A. Person", 91000])
    raw.sheet_state = "hidden"

    ws.row_dimensions[2].hidden = True

    wb.properties.creator = AUTHOR
    wb.properties.lastModifiedBy = "finance-laptop"
    wb.save(dest)

    # openpyxl has no company field; write it into app.xml directly.
    with zipfile.ZipFile(dest) as zin:
        app_xml = set_app_property(zin.read("docProps/app.xml"), "Company", COMPANY)

    # openpyxl will not write a pivot cache or an external link; inject both.
    pivot_definition = (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<pivotCacheDefinition xmlns="http://schemas.openxmlformats.org/'
        b'spreadsheetml/2006/main" recordCount="1"><cacheSource type="worksheet">'
        b'<worksheetSource ref="A1:B2" sheet="RawSalaries"/></cacheSource>'
        b"</pivotCacheDefinition>"
    )
    external_link = (
        b'<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        b'<externalLink xmlns="http://schemas.openxmlformats.org/spreadsheetml/'
        b'2006/main"><externalBook/></externalLink>'
    )
    rewrite_zip(
        dest,
        replacements={"docProps/app.xml": app_xml},
        additions={
            "xl/pivotCache/pivotCacheDefinition1.xml": pivot_definition,
            "xl/externalLinks/externalLink1.xml": external_link,
        },
    )

    return {
        "path": dest.name,
        "author": AUTHOR,
        "company": COMPANY,
        "hidden_sheets": 1,
        "hidden_rows": 1,
        "pivot_caches": 1,
        "external_links": 1,
    }


# ── image ────────────────────────────────────────────────────────────────


def build_photo() -> dict[str, object]:
    from PIL import Image
    from PIL.ExifTags import GPS, Base

    dest = FIXTURE_DIR / "leaky_photo.jpg"
    image = Image.new("RGB", (64, 48), (30, 90, 160))

    exif = image.getexif()
    exif[Base.Make.value] = "Testcorp"
    exif[Base.Model.value] = "Microscope 9000"
    exif[Base.Software.value] = "SecretLabSuite 3.1"
    exif[Base.Artist.value] = AUTHOR
    exif[Base.DateTime.value] = "2026:03:14 15:09:26"

    gps = exif.get_ifd(0x8825)
    gps[GPS.GPSLatitudeRef.value] = "N"
    gps[GPS.GPSLatitude.value] = (39.0, 8.0, 17.48)
    gps[GPS.GPSLongitudeRef.value] = "W"
    gps[GPS.GPSLongitude.value] = (77.0, 13.0, 11.36)

    image.save(dest, exif=exif)
    return {
        "path": dest.name,
        "artist": AUTHOR,
        "has_gps": True,
        "make": "Testcorp",
    }


# ── PDF ──────────────────────────────────────────────────────────────────


def build_pdf() -> dict[str, object]:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.pdfgen import canvas

    dest = FIXTURE_DIR / "leaky_doc.pdf"
    pdf = canvas.Canvas(str(dest), pagesize=LETTER)
    pdf.setAuthor(AUTHOR)
    pdf.setTitle("INTERNAL DRAFT - do not circulate")
    pdf.setSubject("Unblinded results")
    pdf.setCreator("SecretLabSuite 3.1")

    pdf.setFont("Helvetica", 12)
    pdf.drawString(72, 700, "Public summary of the findings.")
    # The classic failed redaction: text first, black box drawn on top.
    pdf.drawString(72, 660, "SECRET: participant P-00417 withdrew after an event.")
    pdf.setFillColorRGB(0, 0, 0)
    pdf.rect(70, 654, 400, 18, stroke=0, fill=1)
    pdf.showPage()
    pdf.save()

    return {
        "path": dest.name,
        "author": AUTHOR,
        "has_possible_redaction": True,
        "hidden_text": "P-00417",
    }


def main() -> None:
    truth = {
        "deck": build_deck(),
        "workbook": build_workbook(),
        "photo": build_photo(),
        "pdf": build_pdf(),
    }
    TRUTH_PATH.write_text(json.dumps(truth, indent=2) + "\n")
    for name, entry in truth.items():
        print(f"{name}: {entry['path']}")
    print(f"truth: {TRUTH_PATH.name}")


if __name__ == "__main__":
    main()
